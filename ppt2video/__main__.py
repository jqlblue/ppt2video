#!/usr/bin/env python

import argparse
import asyncio
import os
import edge_tts
import fitz
from functools import partial
from itertools import count
import locale
from loguru import logger
from operator import is_not, itemgetter
from pathlib import Path
from pptx import Presentation
from pptx.slide import Slide
from string import Template
import subprocess
import tempfile
from typing import Any, TypeGuard
import shutil
import math
from xml.sax.saxutils import escape
from deep_translator import GoogleTranslator
from pypinyin import lazy_pinyin
import re
from num2words import num2words
from pydub import AudioSegment

def get_note_from_slide(slide: Slide) -> str | None:
    if not slide.has_notes_slide:
        return None

    notes_text: str = slide.notes_slide.notes_text_frame.text
    if len(notes_text) == 0:
        return None

    return notes_text

# convert edge tts vtt timestamp to srt format
def ms_to_srt_time(time_unit):
    seconds, milliseconds = divmod(time_unit, 1000)
    minutes, seconds = divmod(seconds, 60)
    hours, minutes = divmod(minutes, 60)
    return f"{hours:02}:{minutes:02}:{seconds:02},{milliseconds:03}"

def replace_pinyin_withtts(pinyin_text):
    start_mapping = dict(
        zh = 'j',
        ch = 'ch',
        z = 'dz',
        c = 'ts',
        # x = 'sh',
        #j = 'dj',
        # q = 'tch',
        # r = 'er'
    )
    for key in start_mapping:
        mapp_item = start_mapping[key]
        if pinyin_text.find(key) != -1:
            return pinyin_text.replace(key, mapp_item)
    return pinyin_text

def replace_pinyin_text(match_text:str, use_tts:bool):
    pinyin_texts = lazy_pinyin(match_text)
    dest_texts = []
    for pinyin_text in pinyin_texts:
        if use_tts:
            pinyin_text = replace_pinyin_withtts(pinyin_text)
        dest_texts.append(f"{pinyin_text}")
    if len(pinyin_texts) > 1:
        target_text =  " ".join(dest_texts)
    else:
        target_text =  "".join(dest_texts)
    return target_text

def replace_markers_target_pinyin_withtts(match):
    match_text = match.group(1)
    target_text = ''
    if match_text.endswith("||"):
        match_texts = match_text.split("||")
        match_text = match_texts[0]
    elif match_text.find("|") != -1:
        match_texts = match_text.split("|")
        target_text = match_texts[1]
        target_text =  replace_pinyin_withtts(target_text)

    if target_text == '':
        target_text =  replace_pinyin_text(match_text, True)
    return target_text

def replace_markers_target_pinyin(match):
    match_text = match.group(1)
    target_text = ''
    if match_text.endswith("||"):
        match_texts = match_text.split("||")
        match_text = match_texts[0]
    elif match_text.find("|") != -1:
        match_texts = match_text.split("|")
        target_text = match_texts[1]

    if target_text == '':
        target_text =  replace_pinyin_text(match_text, False)
    return target_text

def replace_markers_target(match):
    match_text = match.group(1)
    target_text = ''
    if match_text.find("|") != -1:
        match_texts = match_text.split("|")
        target_text = match_texts[1]
    else:
        target_text = match_text

    return target_text

def replace_markers_source_pinyin(match):
    match_text = match.group(1)

    if match_text.endswith("||"):
        match_texts = match_text.split("||")
        target_text = match_texts[0]
        return target_text

    if match_text.find("|") != -1:
        match_texts = match_text.split("|")
        target_text = match_texts[0]
        return target_text

    return match_text

def replace_markers_source(match):
    match_text = match.group(1)

    if match_text.find("|") != -1:
        match_texts = match_text.split("|")
        target_text = match_texts[0]
        return target_text

    return match_text

def replace_numbers(match):
    number = int(match.group())
    return num2words(number)

def get_notes_from_ppt_file(ppt_file_path: Path) -> list[str | None]:
    prs = Presentation(ppt_file_path)
    notes = list(map(get_note_from_slide, prs.slides))
    return notes

async def convert_page_to_image(page: fitz.Page,
                          output_file_path: Path,
                          dpi: int) -> Path:
    loop = asyncio.get_running_loop()
    pix = await loop.run_in_executor(None, partial(page.get_pixmap, dpi=dpi))
    loop.run_in_executor(None, pix.save, output_file_path)
    logger.info('Generate Image file from PDF in `{output_file_path}`', output_file_path=output_file_path)

    return output_file_path

async def convert_ppt_to_image(ppt_file_path: Path,
                               output_dir: Path,
                               dpi: int,
                               output_filename: Template,
                               soffice_file_path: Path,
                               encoding: str,
                               pages: list[int] | None = None) -> list[Path]:

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        output = subprocess.run([soffice_file_path,
                                 '--headless',
                                 '--invisible',
                                 '--convert-to', 'pdf',
                                 '--outdir', tmp_dir_path,
                                 ppt_file_path],
                                check=True, capture_output=True, encoding=encoding)

        if len(output.stderr) > 0:
            raise subprocess.CalledProcessError(output.returncode, output.args, output.stdout, output.stderr)

        pdf_file_path = tmp_dir_path / f'{ppt_file_path.stem}.pdf'
        logger.info('Generate PDF from PPTX in `{pdf_file_path}`', pdf_file_path=pdf_file_path)

        tasks = list()
        with fitz.open(pdf_file_path) as pdf:
            async with asyncio.TaskGroup() as tg:
                for index, page in enumerate(pdf, start=1):
                    if (pages is not None) and (index not in pages):
                        continue

                    output_file_path = output_dir / output_filename.substitute(index=index)
                    tasks.append(tg.create_task(convert_page_to_image(page, output_file_path, dpi)))

    return list(map(lambda t: t.result(), tasks))

async def convert_note_to_audio(note: str,
                          output_file_path: Path,
                          output_subtitles_file_path: Path,
                          souce_lang:str,
                          target_lang:str,
                          voice: str) -> dict:
    show_target_subtitiles = False
    default_chinese_voice = 'zh-CN-XiaoxiaoNeural'

    source_note_paragraphs = []
    target_note_paragraphs = []
    target_voice_paragraphs = []

    notes = note.split("\n")
    notes = list(filter(None, notes))
    notes = [item for item in notes if item.strip()]

    for note_paragraph in notes:
        if souce_lang == target_lang:
            source_note_paragraphs.append(note_paragraph)
            target_note_paragraphs = source_note_paragraphs
            target_voice_paragraphs = source_note_paragraphs
        else:
            show_target_subtitiles = True
            if souce_lang == 'zh-CN':
                replace_markers = re.sub(r'\{(.*?)\}', replace_markers_target, note_paragraph)
                replace_markers_pinyin = re.sub(r'\[(.*?)\]', replace_markers_target_pinyin, replace_markers)
                replace_markers_pinyin_withtts = re.sub(r'\[(.*?)\]', replace_markers_target_pinyin_withtts, replace_markers)
                source_note_pinyin = re.sub(r'\[(.*?)\]', replace_markers_source_pinyin, note_paragraph)
                source_note = re.sub(r'\{(.*?)\}', replace_markers_source, source_note_pinyin)

                source_note_paragraphs.append(source_note)

                if note.find("||") != -1:
                    target_note = replace_markers_pinyin
                    target_voice_note = source_note
                else:
                    target_note = GoogleTranslator(source=souce_lang, target=target_lang).translate(replace_markers_pinyin)
                    target_voice_note = GoogleTranslator(source=souce_lang, target=target_lang).translate(replace_markers_pinyin_withtts)

                target_note_paragraphs.append(target_note)
                target_voice_paragraphs.append(target_voice_note)

    srt_index = 1

    logger.info('Generate Audio Subtitles from target_voice_paragraphs srt_index = `{srt_index}`, `{target_voice_paragraphs}`', target_voice_paragraphs=target_voice_paragraphs,srt_index=srt_index)

    output_acc_paragraphs = []

    start_time = 0
    end_time = 0

    for index, target_voice_paragraph in enumerate(target_voice_paragraphs, start=1):
        output_acc_paragraph = f"{output_file_path}-paragraph-{index}.acc"
        if target_voice_paragraph != '':
            target_note_paragraph = target_note_paragraphs[index-1]
            if show_target_subtitiles and souce_lang == 'zh-CN':
                if note.find("||") != -1:
                    communicate_note_paragraph = edge_tts.Communicate(re.sub(r'\d+', replace_numbers, target_voice_paragraph), default_chinese_voice)
                else:
                    communicate_note_paragraph = edge_tts.Communicate(re.sub(r'\d+', replace_numbers, target_voice_paragraph), voice)
            else:
                communicate_note_paragraph = edge_tts.Communicate(target_note_paragraph, voice)

            await communicate_note_paragraph.save(output_acc_paragraph)

            audio = AudioSegment.from_mp3(output_acc_paragraph)

            duration_ms = len(audio)

            end_time = start_time + duration_ms

            start_time_srt = ms_to_srt_time(start_time)
            end_time_srt = ms_to_srt_time(end_time)

            start_time = end_time

            with open(output_subtitles_file_path, "ab") as srt_file:
                logger.info('Generate Audio Subtitles file from note_paragraph in `{output_subtitles_file_path}`, index = `{index}`', output_subtitles_file_path=output_subtitles_file_path, index = index)

                srt_file.write(f"{srt_index}\n".encode())
                srt_file.write(f"{start_time_srt} --> {end_time_srt}\n".encode())
                srt_file.write(f"{escape(source_note_paragraphs[index-1])}\n\n\n\n".encode())
                if show_target_subtitiles:
                    srt_file.write(f"{escape(target_note_paragraph)}\n\n".encode())

            output_acc_paragraphs.append(output_acc_paragraph)

            # update subtitles index
            srt_index += 1

    combined_audio = AudioSegment.empty()
    for paragraphs_acc in output_acc_paragraphs:
        audio = AudioSegment.from_mp3(paragraphs_acc)
        combined_audio += audio
    combined_audio.export(output_file_path, format="mp3")

    logger.info('Generate Audio file from note in `{output_file_path}`', output_file_path=output_file_path)

    return dict(
        audio=output_file_path,
        subtitles=output_subtitles_file_path
    )

async def convert_notes_to_audio(notes: list[str],
                                 output_dir: Path,
                                 output_filename: Template,
                                 output_subtitles_filename: Template,
                                 source_lang:str,
                                 target_lang:str,
                                 voice: str) -> list[dict]:
    tasks = list()
    async with asyncio.TaskGroup() as tg:
        for index, note in enumerate(notes, start=1):
            output_file_path = output_dir / output_filename.substitute(index=index)
            output_subtitles_file_path = output_dir / output_subtitles_filename.substitute(index=index)
            tasks.append(tg.create_task(convert_note_to_audio(note, output_file_path, output_subtitles_file_path, source_lang, target_lang, voice)))

    return list(map(lambda task: task.result(), tasks))

def convert_video(image_file_path: Path,
                  audio_file_path: Path,
                  audio_subtitles_file_path: Path,
                  output_file_path: Path,
                  ffmpeg_file_path: Path,
                  subtitles_font:str,
                  encoding: str) -> Path:
    output = subprocess.run([ffmpeg_file_path,
                             '-loop', '1',
                             '-i', image_file_path,
                             '-i', audio_file_path,
                             '-vf', f"subtitles={audio_subtitles_file_path}:force_style='BackColour=&HB0000000,Spacing=0.2,Outline=0,Shadow=0.75,Alignment=2,MarginV=25,Fontname={subtitles_font},Fontsize=18,Bold=-1,Borderstyle=3'",
                             '-c:v', 'libx264',
                             '-c:a', 'copy',
                             '-shortest',
                             '-y',
                             output_file_path],
                            check=True, capture_output=True, encoding=encoding)

    logger.info('Generate Video file from Image file and Audio file in `{output_file_path}`', output_file_path=output_file_path)
    return output_file_path

def convert_videos(image_file_paths: list[Path],
                  audio_file_paths: list[dict],
                  output_dir: Path,
                  output_filename: Template,
                  ffmpeg_file_path: Path,
                  subtitles_font:str,
                  encoding: str) -> list[Path]:
    result = list()

    for index, image_file_path, audio_file_path_dict in zip(count(1), image_file_paths, audio_file_paths):
        audio_file_path = audio_file_path_dict['audio']
        audio_subtitles_file_path = audio_file_path_dict['subtitles']
        output_file_path = output_dir / output_filename.substitute(index=index)
        p = convert_video(image_file_path, audio_file_path, audio_subtitles_file_path, output_file_path, ffmpeg_file_path, subtitles_font, encoding)
        result.append(p)

    return result

def concat_videos(video_file_paths: list[Path],
                  output_file_path: Path,
                  ffmpeg_file_path: Path,
                  encoding: str) -> Path:
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        concat_file_path = tmp_dir_path / 'concat.txt'
        with concat_file_path.open(mode='w') as f:
            for p in video_file_paths:
                f.write(f"file '{p.resolve()}'\n")
        output = subprocess.run([ffmpeg_file_path,
                                 '-f', 'concat',
                                 '-safe', '0',
                                 '-i', concat_file_path,
                                 '-c:v', 'copy',
                                 '-c:a', 'aac',
                                 '-ar', '48000',
                                 '-y',
                                 output_file_path],
                                check=True, capture_output=True, encoding=encoding)

        logger.info('Concat Video file from several Video files in `{output_file_path}`', output_file_path=output_file_path)
        return output_file_path

def has_note(t: tuple[int, str | None]) -> TypeGuard[tuple[int, str]]:
    index, note = t
    return (note is not None) and (len(note) > 0)

async def main_process(ppt_file_path: Path,
                       output_file_path: Path,
                       soffice_file_path: Path,
                       ffmpeg_file_path: Path,
                       dpi: int,
                       voice: str,
                       source_lang:str,
                       target_lang:str,
                       subtitles_font:str,
                       encoding: str) -> Path:
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        notes = get_notes_from_ppt_file(ppt_file_path)

        available_pages_and_notes: list[tuple[int, str]] = list(filter(has_note, enumerate(notes, start=1)))
        available_pages: list[int] = list(map(itemgetter(0), available_pages_and_notes))
        available_notes: list[str] = list(map(itemgetter(1), available_pages_and_notes))

        audio_dir_path = tmp_dir_path / 'audios'
        audio_dir_path.mkdir()
        audio_file_paths = await convert_notes_to_audio(notes=available_notes,
                                                        output_dir=audio_dir_path,
                                                        output_filename=Template('note-${index}.aac'),
                                                        output_subtitles_filename=Template('subtitles-${index}.srt'),
                                                        source_lang=source_lang,
                                                        target_lang=target_lang,
                                                        voice=voice)

        image_dir_path = tmp_dir_path / 'images'
        image_dir_path.mkdir()
        image_file_paths = await convert_ppt_to_image(ppt_file_path=ppt_file_path,
                                                      output_dir=image_dir_path,
                                                      pages=available_pages,
                                                      dpi=dpi,
                                                      output_filename=Template('page-${index}.png'),
                                                      soffice_file_path=soffice_file_path,
                                                      encoding=encoding)

        video_dir_path = tmp_dir_path / 'videos'
        video_dir_path.mkdir()
        video_file_paths = convert_videos(image_file_paths=image_file_paths,
                                          audio_file_paths=audio_file_paths,
                                          output_dir=video_dir_path,
                                          output_filename=Template('video-${index}.mp4'),
                                          ffmpeg_file_path=ffmpeg_file_path,
                                          subtitles_font=subtitles_font,
                                          encoding=encoding)

        result = concat_videos(video_file_paths=video_file_paths,
                               output_file_path=output_file_path,
                               ffmpeg_file_path=ffmpeg_file_path,
                               encoding=encoding)

        return result

async def convert(args: argparse.Namespace) -> Path:
    result = await main_process(ppt_file_path=args.infile,
                        output_file_path=args.outfile,
                        soffice_file_path=args.soffice_file_path,
                        ffmpeg_file_path=args.ffmpeg_file_path,
                        dpi=args.dpi,
                        voice=args.voice,
                        source_lang=args.lang,
                        target_lang=args.target_lang,
                        subtitles_font=args.subtitles_font,
                        encoding=args.encoding)
    return result

def pretty_format(obj: dict[Any, Any] | list[Any] | str, depth: int = 0) -> str:
    result = list()

    if isinstance(obj, dict):
        for k, v in obj.items():
            k_str = pretty_format(k, depth)
            if isinstance(v, str):
                result.append(': '.join([k_str, str(v)]))
            else:
                result.append(f'{k_str}:')
                v_str = pretty_format(v, depth+1)
                result.append(v_str)
        if depth == 0:
            result.append('')
    elif isinstance(obj, list):
        for item in obj:
            result.append(pretty_format(item, depth))
    else:
        indent = ' ' * 2 * depth
        result.append(f'{indent}{str(obj)}')

    return '\n'.join(result)

async def list_voices(args: argparse.Namespace) -> None:
    params = dict()

    if args.language != 'all':
        params['Language'] = args.language

    if args.locale != 'all':
        params['Locale'] = args.locale

    if args.gender != 'all':
        params['Gender'] = args.gender.capitalize()

    voices_manager = await edge_tts.VoicesManager.create()
    voices = voices_manager.find(**params)
    voices.sort(key=itemgetter('ShortName'))

    if args.detail == False:
        for name in map(itemgetter('ShortName'), voices):
            print(name)
        return

    print(pretty_format(voices))

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert PPT file(s) to one Video file")
    subparsers = parser.add_subparsers(required=True)

    parser_convert = subparsers.add_parser('convert')
    parser_convert.add_argument('-i', '--infile', type=Path, help='add PPT file(s)', required=True)
    parser_convert.add_argument('outfile', type=Path, help='set output video filename')
    parser_convert.add_argument('--soffice-file-path', type=Path, default=Path(shutil.which('soffice')))
    parser_convert.add_argument('--ffmpeg-file-path', type=Path, default=Path(shutil.which('ffmpeg')))
    parser_convert.add_argument('--dpi', type=int, default=75)
    parser_convert.add_argument('--voice', type=str, default='zh-CN-XiaoxiaoNeural')
    parser_convert.add_argument('--lang', type=str, default='zh-CN')
    parser_convert.add_argument('--target-lang', type=str, default='zh-CN')
    parser_convert.add_argument('--subtitles-font', type=str, default='Arial')
    parser_convert.add_argument('--encoding', type=str, default=locale.getpreferredencoding())
    parser_convert.set_defaults(func=convert)

    parser_list_voices = subparsers.add_parser('list-voices')
    parser_list_voices.add_argument('--language', type=str, default='all')
    parser_list_voices.add_argument('--locale', type=str, default='all')
    parser_list_voices.add_argument('--gender', type=str, default='all')
    parser_list_voices.add_argument('--detail', action='store_true', default=False)
    parser_list_voices.set_defaults(func=list_voices)

    return parser.parse_args()

async def main() -> None:
    args = parse_args()

    await args.func(args)

if __name__ == '__main__':
    asyncio.run(main())
